from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from buho_back.config import settings
import os

templates_directory = settings.TEMPLATES_DIRECTORY


def generate_presentation(content, user_output_files_directory, filename):
    template_path = os.path.join(templates_directory, "powerpoint.pptx")

    prs = Presentation(template_path)

    slide_master = prs.slide_master
    slide_layouts = slide_master.slide_layouts
    slide_width_inches, slide_height_inches = (
        prs.slide_width / 914400,
        prs.slide_height / 914400,
    )

    def create_title_slide(slide_data):
        slide_layout = slide_layouts.get_by_name("Title Slide")
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data["title"]
        slide.placeholders[10].text = slide_data["subtitle"]

    def create_content_slide(slide_data):
        slide_layout = slide_layouts.get_by_name("Title and Content")
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data["title"]
        body = slide.placeholders[13].text_frame
        for point in slide_data["bullet_points"]:
            p = body.add_paragraph()
            p.text = point
            p.level = 0

    def create_table_slide(slide_data):
        slide_layout = slide_layouts.get_by_name("Title Only")
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data["title"]

        # Define table properties
        rows = len(slide_data["table_data"])
        cols = len(slide_data["table_data"][0])
        slide_width_inches
        left = Inches(1)
        top = Inches(2)
        width = Inches(slide_width_inches / 1.18)
        height = Inches(slide_height_inches / 1.56)

        # Add table
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        column_width = table.columns[0].width

        # Write table data
        for row_idx, row_data in enumerate(slide_data["table_data"]):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_data
                # Format the cell
                cell.text_frame.paragraphs[0].font.bold = (
                    True if row_idx == 0 else False
                )
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    RGBColor(0xDD, 0xDD, 0xDD)
                    if row_idx == 0
                    else RGBColor(0xFF, 0xFF, 0xFF)
                )

    # Iterate through the content and add slides accordingly
    for slide_data in content:
        slide_type = slide_data["type"]
        if slide_type == "title":
            create_title_slide(slide_data)
        elif slide_type == "content":
            create_content_slide(slide_data)
        elif slide_type == "table":
            create_table_slide(slide_data)

    # Remove first slide (created by default)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

    if not os.path.exists(user_output_files_directory):
        os.makedirs(user_output_files_directory)

    ppt_file_path = os.path.join(user_output_files_directory, f"{filename}.pptx")
    if os.path.exists(ppt_file_path):
        os.remove(ppt_file_path)
    prs.save(ppt_file_path)

    return ppt_file_path
